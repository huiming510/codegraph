# -*- coding: utf-8 -*-
"""
# @Time    : 2025/11/7 16:56
# @Author  : cuils
# @Description:
"""
import os
import uuid
import base64
from typing import List
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_parser import BaseParser
from ..base import LinkDocument
from ..format_convert import BaseFormatConverter, LibreofficeFormatConverter


class PPTParser(BaseParser):
    def __init__(self, logger, converter: BaseFormatConverter=None):
        super().__init__(logger)
        self.converter = converter
        if self.converter is None:
            self.converter = LibreofficeFormatConverter()

    @property
    def support_file_formats(self) -> List[str]:
        return ["ppt", "pps", "pptx", "ppsx"]

    def parse(self, filepath: str | Path, doc_id=None) -> List[LinkDocument]:
        filepath = Path(filepath).absolute()
        suffix = filepath.suffix.lstrip(".")
        if suffix not in self.support_file_formats:
            raise ValueError(f"Ppt Parser only supports {self.support_file_formats}, but got {suffix}")

        if suffix == "ppt":
            try:
                pptx_file = self.converter.convert(
                    input_filepath=filepath,
                    input_format="ppt",
                    output_format="pptx",
                    output_suffix="pptx"
                )
                filepath = Path(pptx_file).absolute()
            except Exception as e:
                self.logger.info(f"[Error]: 转换失败。`{filepath}` - {repr(e)}")
                return []

        try:
            return self._parse_by_pptx(filepath, doc_id)
        except Exception as e:
            self.logger.info(f"[Parse Error]: 文件解析失败。`{filepath}` - {repr(e)}")
            return []

    def _parse_by_pptx(self, filepath:Path, doc_id=None) -> List[LinkDocument]:
        """使用 pyoffice-pptx解析 pptx格式"""
        ppt = Presentation(filepath)
        markdown_lines = []
        images = []
        for idx, slide in enumerate(ppt.slides, 1):
            slide_lines = [f"## Slide: {idx}\n"]
            shapes = slide.shapes
            for shape in shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text)

                if shape.shape_type in [MSO_SHAPE_TYPE.GROUP]:
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, "text") and sub_shape.text.strip():
                            slide_lines.append(sub_shape.text)

                # 图片、链接图片
                elif shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE]:
                    _byte = shape.image.blob
                    try:
                        _img_type = shape.image.content_type
                    except:
                        _img_type = None
                    _img_name = shape.image.filename
                    _img_ext = _img_name.split(".")[-1].lower()
                    _img_str = base64.b64encode(_byte).decode()
                    md_text = f"\n![media/image]({_img_name})\n"
                    slide_lines.append(md_text)
                    images.append({"img_name": _img_name, "img_ext": _img_ext, "img_str": _img_str})

                # 表格
                elif shape.shape_type in [MSO_SHAPE_TYPE.TABLE]:
                    table_content = []
                    for j, row in enumerate(shape.table.rows):
                        row_content = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_content.append(cell.text)
                        if j == 0:
                            # 表头处理
                            table_content.append(" | "+" | ".join(row_content)+" | ")
                            table_content.append(" | "+" | ".join(["---"]*len(row_content))+" | ")
                        else:
                            table_content.append(" | "+" | ".join(row_content)+" | ")

                    slide_lines.append("")
                    slide_lines.extend(table_content)
                    slide_lines.append("")

                # video/audio
                elif shape.shape_type in [MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.MEDIA]:
                    # TODO：这部分怎么读取、转换成文本
                    pass

                # 嵌入/链接的OLE对象
                elif shape.shape_type in [MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT]:
                    pass

            if slide_lines:
                markdown_lines.extend(slide_lines)
                markdown_lines.append("")

        markdown = "\n".join(markdown_lines)
        document = LinkDocument(
            doc_id=doc_id or str(uuid.uuid4()),
            content=markdown,
            images=images,
            metadata={"file_name": filepath.name}
        )
        return [document]